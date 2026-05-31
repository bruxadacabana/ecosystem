"""
Mnemosyne — Conversor de Slides Markdown (Marp) → .pptx

Recebe o conteúdo Markdown gerado pelo Studio (tipo "Slides") e converte
para um arquivo PowerPoint (.pptx) usando python-pptx.

Formato esperado de entrada (Marp):
  ---
  marp: true
  theme: default
  ---

  # Título do Slide

  Conteúdo do slide aqui
  - Item 1
  - Item 2

  ---

  # Próximo Slide

  Mais conteúdo

  ---

Cada bloco separado por `---` gera um slide. A primeira linha com `#` vira
título; o restante vira corpo de texto com marcadores.
"""
from __future__ import annotations

import logging
import re
from pathlib import Path

log = logging.getLogger("mnemosyne.slidemaker")


def _parse_slides(markdown: str) -> list[dict[str, str]]:
    """Divide o Markdown em slides.

    Retorna lista de dicts {"title": str, "body": str}.
    Slides sem título ficam com title="".
    Frontmatter inicial (entre --- e ---) é ignorado.
    """
    # Remove frontmatter YAML (bloco --- ... --- no topo)
    cleaned = re.sub(r"^---[\s\S]+?---\s*", "", markdown.strip(), count=1)

    # Divide pelos separadores de slide ---
    raw_slides = re.split(r"\n---+\n", cleaned)

    slides: list[dict[str, str]] = []
    for block in raw_slides:
        block = block.strip()
        if not block:
            continue

        lines = block.splitlines()
        title = ""
        body_lines: list[str] = []

        for line in lines:
            # Qualquer nível de heading vira título (pega o primeiro)
            if not title and re.match(r"^#{1,6}\s+", line):
                title = re.sub(r"^#{1,6}\s+", "", line).strip()
            else:
                body_lines.append(line)

        body = "\n".join(body_lines).strip()
        slides.append({"title": title, "body": body})

    return slides


def _add_slide(prs, layout, title: str, body: str) -> None:
    """Adiciona um slide ao prs usando o layout fornecido."""
    from pptx.util import Pt  # type: ignore[import]

    slide = prs.slides.add_slide(layout)

    # Placeholder 0 = título, placeholder 1 = corpo
    if slide.shapes.title:
        slide.shapes.title.text = title or ""

    body_ph = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            body_ph = ph
            break

    if body_ph is None:
        return

    tf = body_ph.text_frame
    tf.word_wrap = True
    tf.clear()

    lines = body.splitlines() if body else []
    first = True
    for line in lines:
        line = line.rstrip()
        if not line:
            continue

        # Detectar marcador de lista: -, *, +, 1.
        is_bullet = bool(re.match(r"^[\-\*\+]\s+|^\d+\.\s+", line))
        text = re.sub(r"^[\-\*\+]\s+|^\d+\.\s+", "", line).strip()

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.text = ("• " if is_bullet else "") + text
        p.font.size = Pt(18)

    # Se não havia linhas, deixar um parágrafo vazio (evita erro de pptx)
    if first:
        tf.paragraphs[0].text = ""


def markdown_to_pptx(
    markdown: str,
    output_path: Path | str,
    title: str = "Apresentação",
) -> Path:
    """Converte Markdown Marp para .pptx.

    Args:
        markdown:    Conteúdo Markdown no formato Marp (gerado pelo Studio).
        output_path: Caminho onde salvar o .pptx (extensão adicionada se ausente).
        title:       Título para a propriedade de documento do .pptx.

    Returns:
        Path do arquivo .pptx gerado.

    Raises:
        ImportError: se python-pptx não estiver instalado.
        ValueError:  se o Markdown não contiver nenhum slide válido.
    """
    try:
        from pptx import Presentation  # type: ignore[import]
        from pptx.util import Inches, Pt  # type: ignore[import]
    except ImportError as exc:
        raise ImportError(
            "python-pptx não encontrado. Instale com: pip install python-pptx"
        ) from exc

    output_path = Path(output_path)
    if output_path.suffix.lower() != ".pptx":
        output_path = output_path.with_suffix(".pptx")

    slides_data = _parse_slides(markdown)
    if not slides_data:
        raise ValueError("Nenhum slide encontrado no Markdown fornecido.")

    log.info("slidemaker: convertendo %d slides para %s", len(slides_data), output_path)

    prs = Presentation()
    prs.core_properties.title = title

    # Layout 1 = "Title and Content" (título + corpo) em todas as themes padrão
    # Layout 0 = "Title Slide" para o primeiro slide
    try:
        title_layout   = prs.slide_layouts[0]  # Title Slide
        content_layout = prs.slide_layouts[1]  # Title and Content
    except IndexError:
        title_layout   = prs.slide_layouts[0]
        content_layout = prs.slide_layouts[0]

    for i, slide_data in enumerate(slides_data):
        layout = title_layout if i == 0 else content_layout
        _add_slide(prs, layout, slide_data["title"], slide_data["body"])
        log.debug("slidemaker: slide %d/%d — '%s'", i + 1, len(slides_data), slide_data["title"])

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    log.info("slidemaker: .pptx salvo em %s (%d slides)", output_path, len(slides_data))
    return output_path
