"""Testes para core/slidemaker.py — parser e conversor Markdown→PPTX."""
from __future__ import annotations

import sys
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest

_ROOT = Path(__file__).parent.parent
if str(_ROOT) not in sys.path:
    sys.path.insert(0, str(_ROOT))

from core.slidemaker import _parse_slides, markdown_to_pptx


# ---------------------------------------------------------------------------
# _parse_slides
# ---------------------------------------------------------------------------

class TestParseSlides:
    def test_basic_two_slides(self):
        md = "# Título A\n\nCorpo A\n\n---\n\n# Título B\n\nCorpo B"
        result = _parse_slides(md)
        assert len(result) == 2
        assert result[0]["title"] == "Título A"
        assert "Corpo A" in result[0]["body"]
        assert result[1]["title"] == "Título B"

    def test_frontmatter_is_stripped(self):
        md = "---\nmarp: true\ntheme: default\n---\n\n# Slide 1\n\nTexto"
        result = _parse_slides(md)
        assert len(result) == 1
        assert result[0]["title"] == "Slide 1"

    def test_slide_without_title(self):
        md = "Apenas corpo sem título"
        result = _parse_slides(md)
        assert len(result) == 1
        assert result[0]["title"] == ""
        assert "Apenas corpo" in result[0]["body"]

    def test_empty_blocks_ignored(self):
        md = "# Slide 1\n\nTexto\n\n---\n\n\n\n---\n\n# Slide 2\n\nTexto 2"
        result = _parse_slides(md)
        # Bloco vazio entre os separadores deve ser ignorado
        assert len(result) == 2

    def test_empty_markdown_returns_empty(self):
        assert _parse_slides("") == []
        assert _parse_slides("   \n\n   ") == []

    def test_heading_levels_all_become_title(self):
        md = "## Subtítulo\n\nCorpo"
        result = _parse_slides(md)
        assert result[0]["title"] == "Subtítulo"

    def test_bullet_lines_in_body(self):
        md = "# Slide\n\n- item 1\n- item 2\n- item 3"
        result = _parse_slides(md)
        assert "item 1" in result[0]["body"]
        assert "item 2" in result[0]["body"]

    def test_only_first_heading_becomes_title(self):
        md = "# Primeiro\n\n## Segundo heading\n\nCorpo"
        result = _parse_slides(md)
        assert result[0]["title"] == "Primeiro"
        assert "Segundo heading" in result[0]["body"]

    def test_many_slides(self):
        slides = "\n\n---\n\n".join(f"# Slide {i}\n\nCorpo {i}" for i in range(8))
        result = _parse_slides(slides)
        assert len(result) == 8

    def test_slide_separator_variations(self):
        """Separadores com mais de 3 hifens devem funcionar."""
        md = "# A\n\nCorpo A\n\n------\n\n# B\n\nCorpo B"
        result = _parse_slides(md)
        assert len(result) == 2


# ---------------------------------------------------------------------------
# markdown_to_pptx
# ---------------------------------------------------------------------------

_MARP_SAMPLE = """\
---
marp: true
---

# Apresentação de Teste

*Subtítulo*

---

## Primeiro Slide

- Ponto 1
- Ponto 2
- Ponto 3

---

## Segundo Slide

Texto livre aqui.

---

## Conclusão

- Takeaway A
- Takeaway B
"""


class TestMarkdownToPptx:
    def test_creates_pptx_file(self, tmp_path):
        out = tmp_path / "test.pptx"
        result = markdown_to_pptx(_MARP_SAMPLE, out)
        assert result == out
        assert out.exists()
        assert out.stat().st_size > 0

    def test_adds_pptx_extension_if_missing(self, tmp_path):
        out = tmp_path / "test"
        result = markdown_to_pptx(_MARP_SAMPLE, out)
        assert result.suffix == ".pptx"
        assert result.exists()

    def test_raises_value_error_on_empty_markdown(self, tmp_path):
        with pytest.raises(ValueError, match="Nenhum slide"):
            markdown_to_pptx("", tmp_path / "out.pptx")

    def test_raises_value_error_on_only_frontmatter(self, tmp_path):
        md = "---\nmarp: true\n---\n\n"
        with pytest.raises(ValueError, match="Nenhum slide"):
            markdown_to_pptx(md, tmp_path / "out.pptx")

    def test_pptx_slide_count_matches_parsed(self, tmp_path):
        from pptx import Presentation  # type: ignore[import]

        out = tmp_path / "count.pptx"
        markdown_to_pptx(_MARP_SAMPLE, out)
        prs = Presentation(str(out))
        expected = len([s for s in _parse_slides(_MARP_SAMPLE) if s["title"] or s["body"]])
        assert len(prs.slides) == expected

    def test_title_property_set(self, tmp_path):
        from pptx import Presentation  # type: ignore[import]

        out = tmp_path / "title.pptx"
        markdown_to_pptx(_MARP_SAMPLE, out, title="Meu Deck")
        prs = Presentation(str(out))
        assert prs.core_properties.title == "Meu Deck"

    def test_creates_parent_dirs(self, tmp_path):
        nested = tmp_path / "a" / "b" / "c" / "out.pptx"
        result = markdown_to_pptx(_MARP_SAMPLE, nested)
        assert result.exists()

    def test_raises_import_error_without_pptx(self, tmp_path, monkeypatch):
        import builtins

        real_import = builtins.__import__

        def mock_import(name, *args, **kwargs):
            if name == "pptx":
                raise ImportError("pptx not available")
            return real_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", mock_import)
        with pytest.raises(ImportError, match="python-pptx"):
            markdown_to_pptx(_MARP_SAMPLE, tmp_path / "out.pptx")

    def test_single_slide_no_separator(self, tmp_path):
        md = "# Título Único\n\n- Ponto A\n- Ponto B"
        out = tmp_path / "single.pptx"
        result = markdown_to_pptx(md, out)
        assert result.exists()

    def test_slide_body_bullet_rendering(self, tmp_path):
        """Bullets devem aparecer com prefixo '• ' no texto do slide."""
        from pptx import Presentation  # type: ignore[import]

        md = "# Slide\n\n- Alpha\n- Beta"
        out = tmp_path / "bullets.pptx"
        markdown_to_pptx(md, out)
        prs = Presentation(str(out))
        all_text = " ".join(
            para.text
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_text_frame
            for para in shape.text_frame.paragraphs
        )
        assert "• Alpha" in all_text
        assert "• Beta" in all_text
